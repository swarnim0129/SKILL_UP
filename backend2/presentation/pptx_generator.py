from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict, Any
import io

# Ported directly from ML_Mumbai with no behavioural changes.

THEME_COLORS = {
    "default": {
        "bg": RGBColor(255, 255, 255),
        "primary": RGBColor(255, 107, 107),
        "secondary": RGBColor(78, 205, 196),
        "accent": RGBColor(255, 230, 109),
        "text": RGBColor(26, 26, 26),
    },
    "modern": {
        "bg": RGBColor(239, 246, 255),
        "primary": RGBColor(99, 102, 241),
        "secondary": RGBColor(139, 92, 246),
        "accent": RGBColor(236, 72, 153),
        "text": RGBColor(30, 41, 59),
    },
    "minimal": {
        "bg": RGBColor(249, 250, 251),
        "primary": RGBColor(0, 0, 0),
        "secondary": RGBColor(64, 64, 64),
        "accent": RGBColor(115, 115, 115),
        "text": RGBColor(0, 0, 0),
    },
    "vibrant": {
        "bg": RGBColor(216, 180, 254),
        "primary": RGBColor(255, 255, 255),
        "secondary": RGBColor(253, 224, 71),
        "accent": RGBColor(52, 211, 153),
        "text": RGBColor(255, 255, 255),
    },
    "dark": {
        "bg": RGBColor(31, 41, 55),
        "primary": RGBColor(96, 165, 250),
        "secondary": RGBColor(52, 211, 153),
        "accent": RGBColor(251, 191, 36),
        "text": RGBColor(249, 250, 251),
    },
}


def create_pptx(
    title: str, slides: List[Dict[str, Any]], theme: str = "default", font: str = "Inter"
) -> io.BytesIO:
    """Generate a PPTX file from slide data."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    colors = THEME_COLORS.get(theme, THEME_COLORS["default"])

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.name = font
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = colors["primary"]

    # Content Slides
    for slide_data in slides:
        layout = slide_data.get("layout", "bullets")
        content = slide_data.get("content", {})

        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["bg"]

        if layout == "bullets":
            _add_bullets_layout(slide, content, colors, font)
        elif layout == "columns":
            _add_columns_layout(slide, content, colors, font)
        elif layout == "timeline":
            _add_timeline_layout(slide, content, colors, font)
        elif layout == "boxes":
            _add_boxes_layout(slide, content, colors, font)
        elif layout == "arrows":
            _add_arrows_layout(slide, content, colors, font)
        elif layout == "compare":
            _add_compare_layout(slide, content, colors, font)
        elif layout == "pyramid":
            _add_pyramid_layout(slide, content, colors, font)
        else:
            _add_default_layout(slide, content, colors, font)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


def _add_bullets_layout(slide, content, colors, font="Inter"):
    """Add bullet point layout with enhanced shapes."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(0.3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["primary"]
    shape.line.color.rgb = colors["primary"]

    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(0.6),
        Inches(9),
        Inches(0.9),
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors["accent"]
    title_bg.line.width = Pt(3)
    title_bg.line.color.rgb = colors["primary"]

    title_box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(0.7),
        Inches(8.6),
        Inches(0.7),
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]

    items = content.get("items", [])
    start_top = 2.0
    for idx, item in enumerate(items):
        bullet_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7),
            Inches(start_top + (idx * 0.65)),
            Inches(0.3),
            Inches(0.3),
        )
        bullet_shape.fill.solid()
        bullet_shape.fill.fore_color.rgb = colors["secondary"]
        bullet_shape.line.width = Pt(0)

        text_box = slide.shapes.add_textbox(
            Inches(1.2),
            Inches(start_top + (idx * 0.65) - 0.05),
            Inches(8.3),
            Inches(0.6),
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = item.get("text", "")
        p.font.name = font
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = colors["text"]

        if item.get("subtext"):
            p_sub = text_frame.add_paragraph()
            p_sub.text = item.get("subtext", "")
            p_sub.font.name = font
            p_sub.font.size = Pt(16)
            p_sub.font.color.rgb = colors["secondary"]
            p_sub.level = 0


def _add_columns_layout(slide, content, colors, font="Inter"):
    """Add two-column layout."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]

    items = content.get("items", [])
    col_width = 4.25
    for idx, item in enumerate(items[:2]):
        left = 0.5 + (idx * 5)
        box = slide.shapes.add_textbox(
            Inches(left), Inches(1.8), Inches(col_width), Inches(3)
        )
        text_frame = box.text_frame

        p_title = text_frame.paragraphs[0]
        p_title.text = item.get("text", "")
        p_title.font.name = font
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = colors["primary"]

        if item.get("subtext"):
            p_content = text_frame.add_paragraph()
            p_content.text = item.get("subtext", "")
            p_content.font.name = font
            p_content.font.size = Pt(16)
            p_content.font.color.rgb = colors["text"]


def _add_timeline_layout(slide, content, colors, font="Inter"):
    """Add timeline layout with enhanced visual elements."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]

    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(2.5),
        Inches(8),
        Inches(0.15),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = colors["secondary"]
    line_shape.line.width = Pt(0)

    items = content.get("items", [])
    if items:
        spacing = 8 / max(len(items), 1)
        for idx, item in enumerate(items):
            left = 1 + (idx * spacing) + (spacing / 2) - 0.4

            shadow_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + 0.05),
                Inches(2.15),
                Inches(0.7),
                Inches(0.7),
            )
            shadow_circle.fill.solid()
            shadow_circle.fill.fore_color.rgb = RGBColor(200, 200, 200)
            shadow_circle.line.width = Pt(0)

            circle_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left),
                Inches(2.1),
                Inches(0.7),
                Inches(0.7),
            )
            circle_box.fill.solid()
            circle_box.fill.fore_color.rgb = colors["primary"]
            circle_box.line.width = Pt(3)
            circle_box.line.color.rgb = RGBColor(255, 255, 255)

            num_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(2.1),
                Inches(0.7),
                Inches(0.7),
            )
            num_frame = num_box.text_frame
            num_frame.text = str(idx + 1)
            num_frame.paragraphs[0].font.name = font
            num_frame.paragraphs[0].font.size = Pt(24)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.vertical_anchor = 1  # Middle

            text_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left - 0.3),
                Inches(3.2),
                Inches(1.3),
                Inches(1.2),
            )
            text_bg.fill.solid()
            text_bg.fill.fore_color.rgb = colors["accent"]
            text_bg.line.width = Pt(2)
            text_bg.line.color.rgb = colors["primary"]

            text_box = slide.shapes.add_textbox(
                Inches(left - 0.2),
                Inches(3.3),
                Inches(1.1),
                Inches(1),
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = item.get("text", "")
            p.font.name = font
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = colors["text"]
            p.alignment = PP_ALIGN.CENTER


def _add_boxes_layout(slide, content, colors, font="Inter"):
    """Add boxes layout with 3D effect."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    items = content.get("items", [])
    if items:
        box_width = 2.6
        spacing = (9 - (len(items) * box_width)) / (len(items) + 1)

        for idx, item in enumerate(items[:3]):
            left = spacing + (idx * (box_width + spacing))

            shadow_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.1),
                Inches(2.1),
                Inches(box_width),
                Inches(2.3),
            )
            shadow_box.fill.solid()
            shadow_box.fill.fore_color.rgb = RGBColor(180, 180, 180)
            shadow_box.line.width = Pt(0)

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(2),
                Inches(box_width),
                Inches(2.3),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = colors["primary"]
            box.line.width = Pt(4)
            box.line.color.rgb = RGBColor(255, 255, 255)

            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + box_width / 2 - 0.3),
                Inches(2.2),
                Inches(0.6),
                Inches(0.6),
            )
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
            icon_circle.line.width = Pt(3)
            icon_circle.line.color.rgb = colors["accent"]

            num_box = slide.shapes.add_textbox(
                Inches(left + box_width / 2 - 0.3),
                Inches(2.2),
                Inches(0.6),
                Inches(0.6),
            )
            num_frame = num_box.text_frame
            num_frame.text = str(idx + 1)
            num_frame.paragraphs[0].font.name = font
            num_frame.paragraphs[0].font.size = Pt(20)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = colors["primary"]
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.vertical_anchor = 1

            text_box = slide.shapes.add_textbox(
                Inches(left + 0.2),
                Inches(3),
                Inches(box_width - 0.4),
                Inches(1.1),
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            p_title = text_frame.paragraphs[0]
            p_title.text = item.get("text", "")
            p_title.font.name = font
            p_title.font.size = Pt(20)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(255, 255, 255)
            p_title.alignment = PP_ALIGN.CENTER

            if item.get("subtext"):
                p_content = text_frame.add_paragraph()
                p_content.text = item.get("subtext", "")
                p_content.font.name = font
                p_content.font.size = Pt(14)
                p_content.font.color.rgb = RGBColor(240, 240, 240)
                p_content.alignment = PP_ALIGN.CENTER


def _add_arrows_layout(slide, content, colors, font="Inter"):
    """Add arrows layout."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]

    items = content.get("items", [])
    if items:
        box_width = 2.5
        arrow_width = 0.5
        total_width = (len(items) * box_width) + ((len(items) - 1) * arrow_width)
        start_left = (10 - total_width) / 2

        for idx, item in enumerate(items):
            left = start_left + (idx * (box_width + arrow_width))
            box = slide.shapes.add_textbox(
                Inches(left),
                Inches(2.5),
                Inches(box_width),
                Inches(1.5),
            )
            text_frame = box.text_frame

            p = text_frame.paragraphs[0]
            p.text = item.get("text", "")
            p.font.name = font
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors["text"]
            p.alignment = PP_ALIGN.CENTER

            if idx < len(items) - 1:
                arrow_box = slide.shapes.add_textbox(
                    Inches(left + box_width),
                    Inches(2.7),
                    Inches(arrow_width),
                    Inches(1),
                )
                arrow_frame = arrow_box.text_frame
                arrow_p = arrow_frame.paragraphs[0]
                arrow_p.text = "→"
                arrow_p.font.name = font
                arrow_p.font.size = Pt(32)
                arrow_p.font.bold = True
                arrow_p.font.color.rgb = colors["primary"]
                arrow_p.alignment = PP_ALIGN.CENTER


def _add_compare_layout(slide, content, colors, font="Inter"):
    """Add comparison layout."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    items = content.get("items", [])
    for idx, item in enumerate(items[:2]):
        left = 0.5 + (idx * 5)
        box = slide.shapes.add_textbox(
            Inches(left),
            Inches(1.8),
            Inches(4.25),
            Inches(3),
        )
        text_frame = box.text_frame

        p_title = text_frame.paragraphs[0]
        p_title.text = item.get("text", "")
        p_title.font.name = font
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = colors["primary"]

        if item.get("subtext"):
            p_content = text_frame.add_paragraph()
            p_content.text = item.get("subtext", "")
            p_content.font.name = font
            p_content.font.size = Pt(16)
            p_content.font.color.rgb = colors["text"]


def _add_pyramid_layout(slide, content, colors, font="Inter"):
    """Add pyramid layout."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]

    items = content.get("items", [])
    if items:
        for idx, item in enumerate(items):
            width = 8 - (idx * 1.5)
            left = (10 - width) / 2
            top = 1.8 + (idx * 0.8)

            box = slide.shapes.add_textbox(
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(0.7),
            )
            text_frame = box.text_frame
            p = text_frame.paragraphs[0]
            p.text = item.get("text", "")
            p.font.name = font
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER


def _add_default_layout(slide, content, colors, font="Inter"):
    """Add default layout."""
    title_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(1.5),
    )
    title_frame = title_box.text_frame
    title_frame.text = content.get("heading", "")
    title_frame.paragraphs[0].font.name = font
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["primary"]
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

